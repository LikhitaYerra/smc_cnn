#!/usr/bin/env python3
"""Generate AI Clinic defense presentation (PowerPoint)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "assets"
FIG = ASSETS / "figures"
OUT = ROOT / "docs" / "AI_Clinic_Defense_Presentation.pptx"

# Brand colours (match digital twin + report)
BG_DARK = RGBColor(0x0A, 0x0F, 0x1A)
BG_PANEL = RGBColor(0x11, 0x18, 0x27)
ACCENT = RGBColor(0x00, 0xD4, 0xAA)
ACCENT_BLUE = RGBColor(0x0B, 0x3D, 0x5C)
WHITE = RGBColor(0xE8, 0xF0, 0xFE)
MUTED = RGBColor(0x88, 0x99, 0xAA)
RED = RGBColor(0xFF, 0x6B, 0x6B)
PURPLE = RGBColor(0x7C, 0x5C, 0xFF)

STUDENTS = [
    "Likhita Yerra",
    "Mohamed Oussama Bouriga",
    "Ahmed Ben Aissa",
    "Abdellahi El Moustapha",
    "Thibault Goutorbe",
]
SUPERVISOR = "Prof. Vishvjit Thakar"
DATE = "31 August 2026"
TITLE = "CNN-Adaptive Sliding Mode Control\nfor Autonomous Robots"
SUBTITLE = "Under Environmental Uncertainty — A Digital Twin Approach"


def _set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_bar(slide, top=Inches(0), height=Inches(0.06), color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(13.333), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _textbox(slide, left, top, width, height, text, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font_name
    p.font.color.rgb = color
    p.alignment = align
    return box


def _bullets(slide, left, top, width, height, items, size=16, color=WHITE,
             bullet_color=ACCENT, spacing=Pt(8)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.level = 0
        p.space_after = spacing
        p.bullet = True
    return box


def _footer(slide, text="PGE5 — AI Clinic  ·  aivancity"):
    _textbox(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.35),
             text, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def _slide_title(slide, title, subtitle=None):
    _set_slide_bg(slide)
    _add_bar(slide, Inches(0.45), Inches(0.05), ACCENT_BLUE)
    _textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7),
             title, size=32, bold=True, color=ACCENT)
    if subtitle:
        _textbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
                 subtitle, size=16, color=MUTED)
    _footer(slide)


def _add_image(slide, path, left, top, width, caption=None):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width)
    if caption:
        _textbox(slide, left, top + width * 0.01 + Inches(0.05), width, Inches(0.4),
                 caption, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def _notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_bar(slide, Inches(0), Inches(0.08), ACCENT)

    logo = ASSETS / "aivancity_logo.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(5.2), Inches(0.35), width=Inches(2.9))

    _textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2),
             TITLE, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.8), Inches(2.65), Inches(11.7), Inches(0.6),
             SUBTITLE, size=18, color=ACCENT, align=PP_ALIGN.CENTER)

    team = "\n".join(STUDENTS)
    _textbox(slide, Inches(2.5), Inches(3.5), Inches(8.3), Inches(1.8),
             f"{team}\n\nSupervisor: {SUPERVISOR}", size=14, color=MUTED,
             align=PP_ALIGN.CENTER)

    _textbox(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
             "PGE5 — AI Clinic  ·  aivancity, Paris  ·  " + DATE,
             size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # Pipeline pill
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(6.2), Inches(8.9), Inches(0.55)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = BG_PANEL
    pill.line.color.rgb = ACCENT
    _textbox(slide, Inches(2.3), Inches(6.25), Inches(8.7), Inches(0.45),
             "Environment  →  AI Agent  →  SMC Control  →  Robot",
             size=13, color=ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "30 sec. Introduce team and supervisor. Mention this is also published at DELCON 2026.")


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Agenda")
    items = [
        "Problem & motivation",
        "CNN-adaptive SMC architecture",
        "Methodology: robot model, CNN, RL agent",
        "Experimental setup & metrics",
        "Results & discussion",
        "3D digital twin demo",
        "Conclusions & future work",
    ]
    _bullets(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(4.5), items, size=20)
    _notes(slide, "45 sec. Total talk ~12 min + 5 min live demo + Q&A.")


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Problem Statement")
    _bullets(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.8), [
        "Autonomous mobile robots must track trajectories under sensor noise, "
        "external disturbances, and wheel slip",
        "Sliding Mode Control (SMC) is robust — but fixed gains cannot optimise "
        "chattering, accuracy, and disturbance rejection simultaneously",
        "Noisy conditions → need smoother control; disturbances/slip → need aggressive correction",
        "Research question: Can a learning-based pipeline adapt SMC parameters "
        "under heterogeneous uncertainty?",
    ], size=17)
    _notes(slide, "1 min. Emphasise the trade-off: one fixed gain cannot do everything.")


def slide_why_smc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Why Adaptive SMC?")
    # Three cards
    cards = [
        ("Classical SMC", "Fixed gains\nStruggles under combined uncertainty", RED),
        ("CNN-Adaptive", "Scenario classification\n→ parameter presets", ACCENT),
        ("RL Agent (PPO)", "Continuous gain adaptation\nfrom reward signal", PURPLE),
    ]
    for i, (title, body, col) in enumerate(cards):
        x = Inches(0.7 + i * 4.1)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0),
                                      Inches(3.7), Inches(3.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_PANEL
        rect.line.color.rgb = col
        _textbox(slide, x + Inches(0.2), Inches(2.2), Inches(3.3), Inches(0.5),
                 title, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
        _textbox(slide, x + Inches(0.2), Inches(2.9), Inches(3.3), Inches(2),
                 body, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    _notes(slide, "45 sec. Three controllers compared in same simulation engine.")


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "System Architecture", "Environment map → CNN → SMC gains → robot")
    _add_image(slide, FIG / "fig01_architecture.png", Inches(0.5), Inches(1.5),
               Inches(12.3), "CNN-adaptive pipeline (top) + PPO RL alternative (bottom)")
    _notes(slide, "1 min. Walk through the pipeline left to right. RL branch is alternative supervisor.")


def slide_robot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Robot & Simulation Model")
    _bullets(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
        "Differential-drive unicycle: state (x, y, θ, v, ω)",
        "Kinematics: ẋ = v cos θ, ẏ = v sin θ, θ̇ = ω",
        "Δt = 0.01 s, wheel base L = 0.3 m",
        "Trajectories: straight, circle, S-curve",
    ], size=16)
    _bullets(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5), [
        "Five uncertainty scenarios:",
        "  Normal — no perturbation",
        "  Noise — Gaussian sensor noise",
        "  Disturbance — impulsive push at t = 8 s",
        "  Slip — 30% velocity reduction t ∈ [10, 14] s",
        "  Combined — all effects together",
    ], size=16)
    _notes(slide, "45 sec. Combined scenario is hardest — used for main demo.")


def slide_cnn(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "CNN Environment Classifier", "EnvironmentCNN — ≈534k parameters")
    _add_image(slide, FIG / "fig02_cnn_samples.png", Inches(0.4), Inches(1.45),
               Inches(7.5))
    _bullets(slide, Inches(8.2), Inches(1.6), Inches(4.5), Inches(5), [
        "Input: 64×64 grayscale occupancy map",
        "Output: 5 scenario classes",
        "3 conv blocks + FC classifier",
        "100% test accuracy (simple maps)",
        "95.1% on realistic maps",
        "Pre-mission classification → SMC preset lookup",
    ], size=15)
    _notes(slide, "1 min. Show the 5 map types. 100% on simple maps validates pipeline.")


def slide_rl(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "RL Agent (PPO)", "Alternative to CNN lookup table")
    _bullets(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5), [
        "Gymnasium environment with domain randomisation",
        "Observation: tracking error, error rate, chattering proxy, scenario flags",
        "Action: select one of 5 SMC parameter presets every 0.2 s",
        "Reward: r = −e − 0.05·j_ω − 0.001·|ω_c|  (error + chattering + effort)",
        "PPO actor-critic with GAE; 80k training steps",
        "Competitive under combined uncertainty; needs more training for disturbance/slip",
    ], size=17)
    _notes(slide, "45 sec. RL is extension work — best chattering under combined, weaker on slip.")


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Evaluation Metrics")
    metrics = [
        ("RMSE tracking error", "Average deviation from reference path"),
        ("Final tracking error", "Error at end of episode — recovery quality"),
        ("Chattering index", "Oscillation in angular velocity command"),
        ("Control effort", "Total |ω_c| — actuator usage"),
    ]
    for i, (name, desc) in enumerate(metrics):
        y = Inches(1.7 + i * 1.25)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8), y, Inches(11.7), Inches(1.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_PANEL
        rect.line.color.rgb = ACCENT_BLUE
        _textbox(slide, Inches(1.0), y + Inches(0.08), Inches(4), Inches(0.4),
                 name, size=18, bold=True, color=ACCENT)
        _textbox(slide, Inches(5.2), y + Inches(0.12), Inches(7), Inches(0.5),
                 desc, size=14, color=MUTED)
    _notes(slide, "30 sec. Chattering = key win for CNN under noise.")


def slide_cnn_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "CNN Classification Results")
    # Big numbers
    for i, (val, label) in enumerate([
        ("100%", "Simple maps\n(test accuracy)"),
        ("95.1%", "Realistic maps\n(with clutter)"),
        ("534k", "Trainable\nparameters"),
    ]):
        x = Inches(1.0 + i * 3.8)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.2), Inches(2.8), Inches(2.8))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BG_PANEL
        circ.line.color.rgb = ACCENT
        _textbox(slide, x, Inches(2.85), Inches(2.8), Inches(0.6),
                 val, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _textbox(slide, x, Inches(3.5), Inches(2.8), Inches(0.8),
                 label, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    _bullets(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5), [
        "Baselines on same split: k-NN 48%, logistic regression 99%, random forest 100%",
        "CNN near oracle ceiling on control benchmarks — interpretable preset selection",
    ], size=14)
    _notes(slide, "45 sec. 95.1% on realistic maps shows deployability beyond toy benchmark.")


def slide_key_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Key Results — CNN-Adaptive vs Classical")
    results = [
        ("35.7%", "Chattering reduction", "under sensor noise"),
        ("14.4%", "Final error improvement", "after disturbance (20.9 → 17.9 mm)"),
        ("32.7%", "Final error improvement", "under wheel slip (40.7 → 27.4 mm)"),
        ("17.9 mm", "CNN final error", "vs oracle 15.6 mm under disturbance"),
    ]
    for i, (num, title, sub) in enumerate(results):
        row, col = divmod(i, 2)
        x = Inches(0.7 + col * 6.2)
        y = Inches(1.7 + row * 2.5)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                                      Inches(5.8), Inches(2.1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_PANEL
        rect.line.color.rgb = ACCENT
        _textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.7),
                 num, size=32, bold=True, color=ACCENT)
        _textbox(slide, x + Inches(0.2), y + Inches(0.85), Inches(5.4), Inches(0.4),
                 title, size=16, bold=True, color=WHITE)
        _textbox(slide, x + Inches(0.2), y + Inches(1.3), Inches(5.4), Inches(0.5),
                 sub, size=12, color=MUTED)
    _notes(slide, "1 min. Hit the three headline numbers: 35.7%, 14.4%, 32.7%. Mention oracle gap only 2.3 mm.")


def slide_figure(prs, title, subtitle, fig_name, caption, img_width=Inches(10.5)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, title, subtitle)
    left = (Inches(13.333) - img_width) / 2
    _add_image(slide, FIG / fig_name, left, Inches(1.55), img_width, caption)
    _notes(slide, f"{title}: point to the figure. {caption or ''}")


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Live Demo Plan", "http://localhost:8000")
    _add_image(slide, FIG / "fig08_digital_twin.png", Inches(0.3), Inches(1.4),
               Inches(7.8))
    _bullets(slide, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5), [
        "Press T — Controller Tour",
        "Classical → CNN → RL auto-run",
        "Compare tab — benchmark all 3",
        "Dual compare — overlay paths",
        "Replay saved runs",
        "Export metrics CSV",
    ], size=14)
    _textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "Keyboard: Space = run/pause  ·  1/2/3 = controllers  ·  R = reset",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    _notes(slide, "Switch to browser NOW. Run: python run_digital_twin.py → localhost:8000. Press T for Controller Tour.")


def slide_conclusions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Conclusions")
    _bullets(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5), [
        "Complete pipeline: dataset → CNN → adaptive SMC → simulation → digital twin",
        "CNN-adaptive SMC reduces chattering and improves recovery — with known trade-offs",
        "Near-oracle performance vs fuzzy and RL baselines on this benchmark",
        "Lyapunov UUB analysis grounds scenario-specific boundary-layer design",
        "3D digital twin enables live defense demonstration and quantitative comparison",
    ], size=17)
    _notes(slide, "45 sec. Summarise contributions — pipeline + theory + twin.")


def slide_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "Future Work")
    _bullets(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(4.5), [
        "Real LiDAR / camera occupancy grids (realistic CNN already at 95.1%)",
        "Continuous on-line scenario inference (mid-mission adaptation)",
        "Extended RL training with curriculum learning & reward shaping",
        "Curved trajectories + Monte-Carlo confidence intervals",
        "Physical differential-drive robot validation",
    ], size=18)
    _notes(slide, "30 sec. Real LiDAR and physical robot are the natural next steps.")


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_bar(slide, Inches(3.2), Inches(0.06), ACCENT)
    _textbox(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1),
             "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6),
             "Questions?", size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.2),
             f"Supervisor: {SUPERVISOR}\n"
             "github.com/LikhitaYerra/smc_cnn\n"
             "aivancity — PGE5 AI Clinic",
             size=16, color=MUTED, align=PP_ALIGN.CENTER)
    logo = ASSETS / "aivancity_logo.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(5.5), Inches(5.8), width=Inches(2.3))
    _notes(slide, "Open floor for questions. Have localhost:8000 ready as backup demo.")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_why_smc(prs)
    slide_architecture(prs)
    slide_robot(prs)
    slide_cnn(prs)
    slide_rl(prs)
    slide_metrics(prs)
    slide_cnn_results(prs)
    slide_key_results(prs)
    slide_figure(prs, "Trajectory Comparison", "Combined scenario — all three controllers",
                 "fig03_trajectory_comparison.png",
                 "Impulsive push visible at x ≈ 2.4 m; adaptive controllers recover tighter")
    slide_figure(prs, "Tracking Error Over Time", "Full 20 s episode under combined uncertainty",
                 "fig04_tracking_error.png",
                 "Disturbance at t = 8 s; slip window t ∈ [10, 14] s")
    slide_figure(prs, "RMSE per Scenario", "Values from published benchmark (Table 3)",
                 "fig05_rmse_bars.png", None, Inches(9))
    slide_figure(prs, "Chattering Index", "Largest CNN gains under noise & combined",
                 "fig06_chattering_bars.png", None, Inches(9))
    slide_figure(prs, "Improvement Heatmap", "CNN-adaptive vs classical SMC (%)",
                 "fig09_improvement_heatmap.png", "Green = adaptive better", Inches(8.5))
    slide_figure(prs, "Multi-Controller Benchmark", "Classical · Fuzzy · CNN · Oracle · RL",
                 "fig10_multicontroller.png", "Lower is better", Inches(9.5))
    slide_demo(prs)
    slide_conclusions(prs)
    slide_future(prs)
    slide_thanks(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    from pptx import Presentation as P
    n = len(P(str(path)).slides)
    print(f"Presentation saved: {path}")
    print(f"Slides: {n}")
