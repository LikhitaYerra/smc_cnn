#!/usr/bin/env python3
"""Generate a polished, editable AI Clinic defense PowerPoint with speaker notes."""

from __future__ import annotations

from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "assets"
FIG = ASSETS / "figures"
OUT = ROOT / "docs" / "AI_Clinic_Defense_Presentation.pptx"
SCRIPT_OUT = ROOT / "docs" / "AI_Clinic_Defense_Speaker_Script.md"

W, H = Inches(13.333), Inches(7.5)

# Academic, aivancity-inspired palette.
NAVY = RGBColor(11, 61, 92)
NAVY_DARK = RGBColor(6, 36, 56)
TEAL = RGBColor(0, 169, 143)
TEAL_LIGHT = RGBColor(220, 246, 241)
GOLD = RGBColor(216, 157, 54)
RED = RGBColor(196, 71, 71)
PURPLE = RGBColor(107, 87, 180)
INK = RGBColor(28, 39, 50)
MUTED = RGBColor(91, 107, 120)
LIGHT = RGBColor(244, 247, 249)
LINE = RGBColor(218, 226, 232)
WHITE = RGBColor(255, 255, 255)

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"

STUDENTS = (
    "Likhita Yerra · Mohamed Oussama Bouriga · Ahmed Ben Aissa · "
    "Abdellahi El Moustapha · Thibault Goutorbe"
)


def set_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, radius=True, line=None, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.fill.transparency = transparency
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def text(
    slide,
    x,
    y,
    w,
    h,
    value,
    size=18,
    color=INK,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def rich_text(slide, x, y, w, h, runs, size=18, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for value, color, bold in runs:
        r = p.add_run()
        r.text = value
        r.font.name = BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return box


def bullets(slide, x, y, w, h, items, size=17, color=INK, spacing=8):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.level = 0
        p.text_frame if False else None
        # Use a colored bullet glyph to avoid theme-dependent bullet formatting.
        p.text = "•  " + p.text
    return box


def add_image_contain(slide, path, x, y, w, h):
    """Fit an image inside a box while preserving its aspect ratio."""
    path = Path(path)
    if not path.exists():
        rect(slide, x, y, w, h, LIGHT, line=LINE)
        text(slide, x, y + h / 2 - Inches(0.2), w, Inches(0.4), "Image unavailable",
             14, MUTED, align=PP_ALIGN.CENTER)
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    image_ratio = iw / ih
    if image_ratio > box_ratio:
        pw = w
        ph = w / image_ratio
    else:
        ph = h
        pw = h * image_ratio
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)


def header(slide, title_value, section, slide_no):
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), W, Inches(0.16), TEAL, radius=False)
    text(slide, Inches(0.62), Inches(0.37), Inches(11.9), Inches(0.45),
         title_value, 26, NAVY_DARK, True, TITLE_FONT)
    text(slide, Inches(0.65), Inches(0.92), Inches(5.0), Inches(0.25),
         section.upper(), 9, TEAL, True)
    rect(slide, Inches(0.62), Inches(1.2), Inches(12.05), Inches(0.02), LINE, radius=False)
    footer(slide, slide_no)


def footer(slide, slide_no):
    text(slide, Inches(0.62), Inches(7.08), Inches(8.5), Inches(0.22),
         "PGE5 · AI Clinic · aivancity", 9, MUTED)
    text(slide, Inches(11.7), Inches(7.08), Inches(0.95), Inches(0.22),
         f"{slide_no:02d}", 9, NAVY, True, align=PP_ALIGN.RIGHT)


def note(slide, script):
    notes = slide.notes_slide.notes_text_frame
    notes.text = script


def takeaway(slide, value, y=Inches(6.38), color=TEAL_LIGHT):
    rect(slide, Inches(0.62), y, Inches(12.05), Inches(0.52), color, line=TEAL)
    text(slide, Inches(0.84), y + Inches(0.11), Inches(11.6), Inches(0.28),
         "TAKEAWAY  " + value, 12, NAVY_DARK, True)


def metric_card(slide, x, y, w, number, title_value, detail, color=TEAL):
    rect(slide, x, y, w, Inches(1.48), WHITE, line=LINE)
    rect(slide, x, y, Inches(0.09), Inches(1.48), color, radius=False)
    text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
         number, 28, color, True, TITLE_FONT)
    text(slide, x + Inches(0.25), y + Inches(0.69), w - Inches(0.4), Inches(0.28),
         title_value, 13, INK, True)
    text(slide, x + Inches(0.25), y + Inches(1.03), w - Inches(0.4), Inches(0.3),
         detail, 10, MUTED)


def process_node(slide, x, y, w, title_value, subtitle, color):
    rect(slide, x, y, w, Inches(1.05), WHITE, line=color)
    text(slide, x + Inches(0.1), y + Inches(0.17), w - Inches(0.2), Inches(0.3),
         title_value, 15, color, True, align=PP_ALIGN.CENTER)
    text(slide, x + Inches(0.1), y + Inches(0.56), w - Inches(0.2), Inches(0.25),
         subtitle, 9, MUTED, align=PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=TEAL):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]
    scripts = []

    # 1 — Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, LIGHT)
    rect(slide, Inches(0), Inches(0), Inches(7.35), H, NAVY_DARK, radius=False)
    rect(slide, Inches(0.65), Inches(0.62), Inches(1.15), Inches(0.08), TEAL, radius=False)
    text(slide, Inches(0.65), Inches(0.9), Inches(6.1), Inches(1.65),
         "CNN-Adaptive\nSliding Mode Control", 34, WHITE, True, TITLE_FONT)
    text(slide, Inches(0.68), Inches(2.66), Inches(5.95), Inches(0.7),
         "Autonomous differential-drive robots under environmental uncertainty",
         17, RGBColor(200, 220, 230))
    text(slide, Inches(0.68), Inches(3.58), Inches(5.9), Inches(0.38),
         "A DIGITAL TWIN APPROACH", 11, TEAL, True)
    text(slide, Inches(0.68), Inches(4.33), Inches(5.95), Inches(0.85),
         STUDENTS, 12, WHITE)
    text(slide, Inches(0.68), Inches(5.42), Inches(5.95), Inches(0.5),
         "Supervisor  ·  Prof. Vishvjit Thakar", 13, RGBColor(200, 220, 230), True)
    text(slide, Inches(0.68), Inches(6.25), Inches(5.95), Inches(0.35),
         "PGE5 · AI Clinic · 31 August 2026", 11, WHITE)
    add_image_contain(slide, FIG / "fig08_digital_twin.png",
                      Inches(7.35), Inches(0.65), Inches(5.65), Inches(5.55))
    rect(slide, Inches(7.35), Inches(6.25), Inches(5.65), Inches(0.6), WHITE, line=LINE)
    add_image_contain(slide, ASSETS / "aivancity_logo.png",
                      Inches(8.7), Inches(6.32), Inches(2.9), Inches(0.42))
    script = (
        "[Timing: 40 seconds]\n\n"
        "Good morning. We are presenting our AI Clinic project on CNN-adaptive sliding "
        "mode control for autonomous differential-drive robots under environmental "
        "uncertainty. Our goal was not only to build a robust controller, but to create "
        "a complete and explainable AI-control pipeline—from environment perception to "
        "controller adaptation and a live 3D digital twin. The project was completed by "
        "our five-person team under the supervision of Professor Vishvjit Thakar. I will "
        "first explain the control problem, then the CNN and PPO approaches, show the "
        "quantitative results, and finish with the digital twin demonstration."
    )
    note(slide, script); scripts.append(("1. Title", script))

    # 2 — Story map
    slide = prs.slides.add_slide(blank); header(slide, "The Defense in One Story", "Roadmap", 2)
    stages = [
        ("01", "Problem", "Fixed SMC gains fail to balance accuracy and smoothness"),
        ("02", "Method", "CNN recognises context and selects interpretable gains"),
        ("03", "Evidence", "Five scenarios, four metrics, multiple baselines"),
        ("04", "Impact", "Live digital twin makes performance visible and repeatable"),
    ]
    for i, (n, t, d) in enumerate(stages):
        x = Inches(0.72 + i * 3.05)
        text(slide, x, Inches(1.7), Inches(0.6), Inches(0.45), n, 23, TEAL, True, TITLE_FONT)
        rect(slide, x, Inches(2.28), Inches(2.65), Inches(2.45), WHITE, line=LINE)
        text(slide, x + Inches(0.2), Inches(2.56), Inches(2.25), Inches(0.42), t, 20, NAVY, True)
        text(slide, x + Inches(0.2), Inches(3.15), Inches(2.25), Inches(1.1), d, 13, MUTED)
        if i < 3:
            arrow(slide, x + Inches(2.7), Inches(3.48), x + Inches(2.98), Inches(3.48))
    takeaway(slide, "The project connects robust control, explainable AI and an interactive digital twin.")
    script = (
        "[Timing: 35 seconds]\n\n"
        "The presentation follows one simple story. First, a fixed-gain controller cannot "
        "respond optimally to every operating condition. Second, we use a CNN to identify "
        "the condition and choose an interpretable controller preset. Third, we evaluate "
        "that idea consistently across five scenarios and against several baselines. "
        "Finally, we expose the entire system through a live digital twin. This roadmap is "
        "important because the digital twin is not a separate interface; it is the final "
        "layer of the same research pipeline."
    )
    note(slide, script); scripts.append(("2. Defense story", script))

    # 3 — Problem
    slide = prs.slides.add_slide(blank); header(slide, "One Controller, Conflicting Requirements", "Problem", 3)
    rect(slide, Inches(0.7), Inches(1.55), Inches(3.7), Inches(4.55), LIGHT, line=LINE)
    text(slide, Inches(0.98), Inches(1.87), Inches(3.1), Inches(0.45), "Fixed-gain SMC", 23, NAVY, True, TITLE_FONT)
    bullets(slide, Inches(1.0), Inches(2.55), Inches(3.0), Inches(2.8), [
        "Strong robustness to matched disturbances",
        "Simple closed-form control law",
        "Same gains in every environment",
        "Chattering–accuracy trade-off",
    ], 15)
    conditions = [
        ("Sensor noise", "needs wider boundary layer\nand stronger smoothing", TEAL),
        ("External push", "needs faster heading recovery\nand higher correction gain", RED),
        ("Wheel slip", "needs larger surface gains\nto compensate velocity loss", GOLD),
    ]
    for i, (t, d, c) in enumerate(conditions):
        y = Inches(1.55 + i * 1.53)
        rect(slide, Inches(4.85), y, Inches(7.75), Inches(1.2), WHITE, line=c)
        text(slide, Inches(5.18), y + Inches(0.18), Inches(2.15), Inches(0.35), t, 17, c, True)
        text(slide, Inches(7.45), y + Inches(0.18), Inches(4.7), Inches(0.62), d, 13, MUTED)
    takeaway(slide, "A single static parameter set cannot optimise all operating regimes.")
    script = (
        "[Timing: 60 seconds]\n\n"
        "Sliding mode control is a strong baseline because it is robust to disturbances "
        "and modelling errors. The limitation is that its behavior depends directly on "
        "fixed gains and the boundary-layer width. Sensor noise requires more smoothing "
        "to avoid rapid actuator switching. An external push requires more aggressive "
        "heading correction. Wheel slip requires stronger gains to recover the lost "
        "motion. These requirements conflict. If we tune only for noise, recovery becomes "
        "slower; if we tune only for disturbance rejection, chattering can increase. "
        "This is the central problem: one static parameter set cannot be optimal in all "
        "five environments."
    )
    note(slide, script); scripts.append(("3. Problem statement", script))

    # 4 — Gap and contributions
    slide = prs.slides.add_slide(blank); header(slide, "Research Gap & Contributions", "Introduction", 4)
    rect(slide, Inches(0.7), Inches(1.55), Inches(5.05), Inches(4.65), NAVY_DARK)
    text(slide, Inches(1.0), Inches(1.88), Inches(4.45), Inches(0.45), "Research gap", 23, WHITE, True, TITLE_FONT)
    text(slide, Inches(1.0), Inches(2.6), Inches(4.4), Inches(2.6),
         "SMC, CNN perception and reinforcement learning are well studied individually. "
         "Few compact, reproducible systems connect environment classification directly "
         "to interpretable SMC gain scheduling and evaluate the full chain in a digital twin.",
         16, RGBColor(220, 232, 238))
    contributions = [
        ("1", "Reproducible pipeline", "Map generation → CNN → controller → evaluation"),
        ("2", "Explainable adaptation", "Five classes mapped to hand-tuned gain presets"),
        ("3", "Multi-controller benchmark", "Classical, fuzzy, CNN, oracle and PPO"),
        ("4", "Defense-ready digital twin", "Live simulation, replay, compare and export"),
    ]
    for i, (n, t, d) in enumerate(contributions):
        y = Inches(1.55 + i * 1.13)
        text(slide, Inches(6.15), y + Inches(0.06), Inches(0.42), Inches(0.4), n, 18, TEAL, True, TITLE_FONT)
        text(slide, Inches(6.72), y, Inches(5.45), Inches(0.32), t, 15, NAVY, True)
        text(slide, Inches(6.72), y + Inches(0.4), Inches(5.45), Inches(0.38), d, 11, MUTED)
    takeaway(slide, "The novelty is the integrated, interpretable and demonstrable control pipeline.")
    script = (
        "[Timing: 55 seconds]\n\n"
        "The individual technologies are not new by themselves. Sliding mode control, "
        "CNN perception and reinforcement learning each have a large literature. Our "
        "research gap is the integration of these pieces into one compact and reproducible "
        "system where the CNN does not replace the controller—it supervises it. That keeps "
        "the actual control law interpretable. Our four contributions are the full pipeline, "
        "scenario-aware gain presets, a multi-controller benchmark including an oracle "
        "upper bound, and a live digital twin that makes the research repeatable during "
        "the defense."
    )
    note(slide, script); scripts.append(("4. Gap and contributions", script))

    # 5 — Pipeline
    slide = prs.slides.add_slide(blank); header(slide, "From Perception to Control", "Architecture", 5)
    nodes = [
        ("Environment map", "64×64 occupancy map", NAVY),
        ("CNN classifier", "5 scenario classes", TEAL),
        ("Preset selector", "λ, k, φ, αω", GOLD),
        ("SMC controller", "v and ω commands", PURPLE),
        ("Robot", "differential drive", RED),
    ]
    for i, (t, d, c) in enumerate(nodes):
        x = Inches(0.55 + i * 2.55)
        process_node(slide, x, Inches(2.15), Inches(2.0), t, d, c)
        if i < 4:
            arrow(slide, x + Inches(2.0), Inches(2.68), x + Inches(2.48), Inches(2.68))
    rect(slide, Inches(3.04), Inches(4.05), Inches(7.25), Inches(1.2), TEAL_LIGHT, line=TEAL)
    text(slide, Inches(3.25), Inches(4.27), Inches(6.8), Inches(0.35),
         "Alternative branch: PPO agent", 18, NAVY, True, align=PP_ALIGN.CENTER)
    text(slide, Inches(3.25), Inches(4.72), Inches(6.8), Inches(0.28),
         "Robot state + error → policy → adaptive SMC parameters", 12, MUTED, align=PP_ALIGN.CENTER)
    takeaway(slide, "The CNN changes controller parameters—not the safety-critical control structure.")
    script = (
        "[Timing: 60 seconds]\n\n"
        "This is the main system architecture. A grayscale occupancy map represents the "
        "operating condition. The lightweight CNN predicts one of five classes. That class "
        "selects a scenario-specific parameter preset containing the sliding-surface gains, "
        "velocity gains, boundary-layer width and angular smoothing. The SMC then calculates "
        "linear and angular velocity commands exactly as before. This distinction is important: "
        "the learned model is a supervisor, not an opaque end-to-end controller. The bottom "
        "branch shows our PPO alternative, which uses robot state and error to adapt parameters."
    )
    note(slide, script); scripts.append(("5. Architecture", script))

    # 6 — Model and control law
    slide = prs.slides.add_slide(blank); header(slide, "Robot Model & Control Law", "Methodology", 6)
    rect(slide, Inches(0.7), Inches(1.52), Inches(5.7), Inches(4.75), LIGHT, line=LINE)
    text(slide, Inches(1.0), Inches(1.85), Inches(5.1), Inches(0.38),
         "Differential-drive kinematics", 20, NAVY, True, TITLE_FONT)
    text(slide, Inches(1.05), Inches(2.55), Inches(5.0), Inches(1.25),
         "ẋ = v cos θ\nẏ = v sin θ\nθ̇ = ω", 24, NAVY_DARK, True, TITLE_FONT, align=PP_ALIGN.CENTER)
    text(slide, Inches(1.0), Inches(4.2), Inches(5.1), Inches(1.0),
         "State q = [x, y, θ, v, ω]ᵀ\nΔt = 0.01 s  ·  wheel base L = 0.3 m",
         14, MUTED, align=PP_ALIGN.CENTER)
    rect(slide, Inches(6.75), Inches(1.52), Inches(5.88), Inches(4.75), WHITE, line=LINE)
    text(slide, Inches(7.08), Inches(1.85), Inches(5.25), Inches(0.38),
         "Sliding-mode control", 20, NAVY, True, TITLE_FONT)
    formulas = [
        ("Sliding surface", "sᵢ = ėᵢ + λᵢeᵢ"),
        ("Linear command", "vᶜ = vʳ + kᵥdᶠ"),
        ("Angular switching", "ωraw = kω[ψ + tanh(sᵧ/φ) + 0.2 tanh(sθ/φ)]"),
        ("Smoothing", "ωᶜ = αωωᶜₖ₋₁ + (1−αω)ωraw"),
    ]
    for i, (label, formula) in enumerate(formulas):
        y = Inches(2.5 + i * 0.78)
        text(slide, Inches(7.05), y, Inches(1.55), Inches(0.28), label, 11, TEAL, True)
        text(slide, Inches(8.65), y - Inches(0.02), Inches(3.6), Inches(0.36), formula, 13, INK, False, TITLE_FONT)
    takeaway(slide, "The boundary layer φ directly controls the robustness–chattering trade-off.")
    script = (
        "[Timing: 70 seconds]\n\n"
        "The robot is modeled as a standard differential-drive unicycle. Its position evolves "
        "through linear speed v and heading through angular speed omega. The non-holonomic "
        "constraint means it cannot move sideways, so lateral pushes are challenging. On the "
        "control side, we define sliding surfaces from tracking error and its derivative. "
        "The linear command corrects forward error. The angular command combines bearing "
        "correction with smooth hyperbolic-tangent switching. Finally, exponential smoothing "
        "reduces high-frequency oscillation. The most important design variable is phi: a "
        "wider boundary layer reduces chattering but permits more steady-state error."
    )
    note(slide, script); scripts.append(("6. Robot and control law", script))

    # 7 — Scenarios
    slide = prs.slides.add_slide(blank); header(slide, "Five Controlled Test Scenarios", "Experimental design", 7)
    scenarios = [
        ("Normal", "No perturbation", NAVY),
        ("Noise", "Gaussian pose noise\nσxy = 0.02 m", TEAL),
        ("Disturbance", "Lateral push +0.4 m\nat t = 8 s", RED),
        ("Slip", "30% velocity reduction\nt = 10–14 s", GOLD),
        ("Combined", "Noise + push + slip", PURPLE),
    ]
    for i, (t, d, c) in enumerate(scenarios):
        x = Inches(0.48 + i * 2.55)
        rect(slide, x, Inches(1.72), Inches(2.25), Inches(3.55), WHITE, line=c)
        text(slide, x + Inches(0.18), Inches(2.03), Inches(1.9), Inches(0.4),
             f"{i+1:02d}", 19, c, True, TITLE_FONT)
        text(slide, x + Inches(0.18), Inches(2.65), Inches(1.9), Inches(0.4),
             t, 18, NAVY_DARK, True)
        text(slide, x + Inches(0.18), Inches(3.36), Inches(1.9), Inches(0.8),
             d, 13, MUTED, align=PP_ALIGN.CENTER)
    takeaway(slide, "All controllers are evaluated under identical perturbations and timing.")
    script = (
        "[Timing: 45 seconds]\n\n"
        "We use five controlled scenarios. Normal is the sanity check. Noise corrupts pose "
        "measurements. Disturbance applies a lateral push at eight seconds. Slip reduces "
        "effective velocity by thirty percent between ten and fourteen seconds. Combined "
        "applies all effects together and is therefore the hardest case. Every controller "
        "uses the same robot model, reference speed, simulation step and perturbation timing, "
        "so the comparison is fair. The normal scenario also verifies that adaptation does "
        "not introduce unnecessary behavior when no uncertainty exists."
    )
    note(slide, script); scripts.append(("7. Scenarios", script))

    # 8 — CNN
    slide = prs.slides.add_slide(blank); header(slide, "EnvironmentCNN: Lightweight Scenario Recognition", "AI Method", 8)
    add_image_contain(slide, FIG / "fig02_cnn_samples.png",
                      Inches(0.65), Inches(1.55), Inches(7.05), Inches(3.3))
    rect(slide, Inches(8.02), Inches(1.55), Inches(4.6), Inches(4.55), LIGHT, line=LINE)
    text(slide, Inches(8.35), Inches(1.86), Inches(4.0), Inches(0.4),
         "Architecture", 20, NAVY, True, TITLE_FONT)
    bullets(slide, Inches(8.35), Inches(2.45), Inches(3.85), Inches(2.15), [
        "Input: 1 × 64 × 64",
        "Conv blocks: 16 → 32 → 64",
        "BatchNorm + ReLU + MaxPool",
        "FC 4096 → 128 → 5",
        "≈534k trainable parameters",
    ], 13, spacing=5)
    metric_card(slide, Inches(0.75), Inches(5.02), Inches(2.15), "100%", "Simple maps", "225 held-out test maps", TEAL)
    metric_card(slide, Inches(3.05), Inches(5.02), Inches(2.15), "95.1%", "Realistic maps", "walls, clutter & obstacles", NAVY)
    metric_card(slide, Inches(5.35), Inches(5.02), Inches(2.15), "5", "Scenario classes", "normal to combined", GOLD)
    takeaway(slide, "Classification is fast, interpretable and performed once before the mission.", Inches(6.58))
    script = (
        "[Timing: 65 seconds]\n\n"
        "The CNN receives a sixty-four by sixty-four grayscale map. Each scenario has a "
        "distinct visual signature: a clean path, noisy pixels, an impact marker, a slip "
        "band, or all cues combined. The network has three convolutional blocks followed "
        "by a small classifier, for approximately five hundred and thirty-four thousand "
        "parameters. It reaches one hundred percent on the engineered test set. Because "
        "that dataset is deliberately separable, we also created more realistic cluttered "
        "maps; the CNN still reaches ninety-five point one percent. In deployment, one "
        "pre-mission inference selects the controller preset, so runtime cost is minimal."
    )
    note(slide, script); scripts.append(("8. CNN", script))

    # 9 — Gain logic
    slide = prs.slides.add_slide(blank); header(slide, "How the Presets Change Controller Behaviour", "Adaptive control", 9)
    rows = [
        ("Normal", "nominal gains", "baseline behaviour", NAVY),
        ("Noise", "↑ φ, ↑ smoothing", "less angular chattering", TEAL),
        ("Disturbance", "↑ λy, ↑ kω", "faster heading recovery", RED),
        ("Slip", "↑ λx, ↑ λy", "compensate velocity deficit", GOLD),
        ("Combined", "balanced preset", "smoothness + correction", PURPLE),
    ]
    text(slide, Inches(0.9), Inches(1.52), Inches(2.5), Inches(0.3), "SCENARIO", 10, MUTED, True)
    text(slide, Inches(3.55), Inches(1.52), Inches(3.0), Inches(0.3), "PARAMETER LOGIC", 10, MUTED, True)
    text(slide, Inches(7.2), Inches(1.52), Inches(4.5), Inches(0.3), "CONTROL EFFECT", 10, MUTED, True)
    for i, (s, logic, effect, c) in enumerate(rows):
        y = Inches(1.92 + i * 0.82)
        rect(slide, Inches(0.75), y, Inches(11.85), Inches(0.64), LIGHT if i % 2 == 0 else WHITE)
        rect(slide, Inches(0.75), y, Inches(0.07), Inches(0.64), c, radius=False)
        text(slide, Inches(0.98), y + Inches(0.14), Inches(2.0), Inches(0.28), s, 13, c, True)
        text(slide, Inches(3.52), y + Inches(0.14), Inches(3.0), Inches(0.28), logic, 13, INK, True)
        text(slide, Inches(7.18), y + Inches(0.14), Inches(4.6), Inches(0.28), effect, 13, MUTED)
    takeaway(slide, "Adaptation remains explainable: each prediction maps to a documented engineering choice.")
    script = (
        "[Timing: 55 seconds]\n\n"
        "This slide explains the system's interpretability. For each class we know exactly "
        "which parameters change and why. Under noise we widen phi and increase smoothing, "
        "which suppresses angular chattering. Under disturbance we increase lateral surface "
        "and angular correction gains, which improves recovery after the push. Under slip "
        "we raise both positional sliding gains to compensate for lost velocity. The combined "
        "preset balances these competing requirements. This is a key advantage over an "
        "end-to-end neural controller: every adaptation can be inspected, justified and "
        "bounded."
    )
    note(slide, script); scripts.append(("9. Preset logic", script))

    # 10 — PPO
    slide = prs.slides.add_slide(blank); header(slide, "PPO Agent: A Continuous-Learning Alternative", "Reinforcement learning", 10)
    flow = [
        ("Observation", "error, rate,\njitter, flags", NAVY),
        ("Actor–critic", "PPO + GAE", PURPLE),
        ("Action", "select SMC\npreset", GOLD),
        ("Reward", "−error − jitter\n− effort", TEAL),
    ]
    for i, (t, d, c) in enumerate(flow):
        x = Inches(0.8 + i * 3.05)
        process_node(slide, x, Inches(1.75), Inches(2.25), t, d, c)
        if i < 3:
            arrow(slide, x + Inches(2.25), Inches(2.28), x + Inches(2.92), Inches(2.28), PURPLE)
    rect(slide, Inches(0.8), Inches(3.55), Inches(11.5), Inches(2.1), LIGHT, line=LINE)
    metric_card(slide, Inches(1.05), Inches(3.85), Inches(2.6), "80k", "Training steps", "40 randomised episodes", PURPLE)
    metric_card(slide, Inches(3.95), Inches(3.85), Inches(2.6), "0.2 s", "Update interval", "policy selects preset", GOLD)
    metric_card(slide, Inches(6.85), Inches(3.85), Inches(2.6), "48.9", "Combined chatter", "best in Table V", TEAL)
    text(slide, Inches(9.7), Inches(4.08), Inches(2.3), Inches(0.9),
         "Promising, but weaker under disturbance and slip → more training needed",
         13, RED, True, align=PP_ALIGN.CENTER)
    takeaway(slide, "PPO is competitive under combined uncertainty but does not yet replace CNN supervision.")
    script = (
        "[Timing: 60 seconds]\n\n"
        "We also implemented PPO as a learning-based comparison. The observation contains "
        "tracking error, error rate, a smoothed jitter measure and uncertainty flags. The "
        "actor-critic policy selects one of the same five SMC presets every zero point two "
        "seconds. Its reward penalizes error, chattering and control effort. After eighty "
        "thousand training steps, PPO reaches the lowest combined-scenario chattering value "
        "of forty-eight point nine. However, it performs worse under disturbance and slip. "
        "Therefore, our conclusion is balanced: PPO is promising, but the CNN supervisor is "
        "currently more reliable and more interpretable for this dataset."
    )
    note(slide, script); scripts.append(("10. PPO", script))

    # 11 — Protocol
    slide = prs.slides.add_slide(blank); header(slide, "Evaluation Protocol", "Experimental setup", 11)
    left_items = [
        ("Duration", "20 s · 2,000 steps"),
        ("Reference", "straight path · 0.3 m/s"),
        ("Randomness", "fixed seeds for reproducibility"),
        ("Controllers", "Classical · Fuzzy · CNN · Oracle · RL"),
    ]
    for i, (k, v) in enumerate(left_items):
        y = Inches(1.6 + i * 1.02)
        text(slide, Inches(0.85), y, Inches(1.4), Inches(0.32), k.upper(), 10, TEAL, True)
        text(slide, Inches(2.2), y - Inches(0.03), Inches(3.6), Inches(0.4), v, 16, INK, True)
        rect(slide, Inches(0.85), y + Inches(0.55), Inches(4.95), Inches(0.015), LINE, radius=False)
    metrics = [
        ("RMSE", "overall tracking accuracy"),
        ("Final error", "post-perturbation recovery"),
        ("Chattering", "command smoothness"),
        ("Control effort", "actuator demand"),
    ]
    rect(slide, Inches(6.35), Inches(1.55), Inches(6.0), Inches(4.55), NAVY_DARK)
    text(slide, Inches(6.72), Inches(1.88), Inches(5.3), Inches(0.4), "Four evaluation metrics", 20, WHITE, True, TITLE_FONT)
    for i, (k, v) in enumerate(metrics):
        y = Inches(2.6 + i * 0.72)
        text(slide, Inches(6.78), y, Inches(1.45), Inches(0.3), k, 14, TEAL, True)
        text(slide, Inches(8.4), y, Inches(3.4), Inches(0.3), v, 13, RGBColor(220, 232, 238))
    takeaway(slide, "Lower error is not enough: a useful controller must also be smooth and efficient.")
    script = (
        "[Timing: 45 seconds]\n\n"
        "The main experiment lasts twenty seconds at a ten-millisecond time step. The robot "
        "tracks a straight reference at zero point three meters per second. Fixed seeds make "
        "the perturbations reproducible. We compare five controller variants, including an "
        "oracle that knows the true scenario and therefore defines the practical ceiling. "
        "We report four metrics because tracking error alone is incomplete. RMSE captures "
        "overall accuracy, final error measures recovery, chattering reflects actuator "
        "smoothness, and control effort reflects energy and mechanical demand."
    )
    note(slide, script); scripts.append(("11. Evaluation protocol", script))

    # 12 — Results headline
    slide = prs.slides.add_slide(blank); header(slide, "The Four Numbers to Remember", "Results", 12)
    metric_card(slide, Inches(0.72), Inches(1.65), Inches(2.85), "35.7%", "Less chattering", "noise scenario · 76.48 → 49.14", TEAL)
    metric_card(slide, Inches(3.72), Inches(1.65), Inches(2.85), "14.4%", "Better final error", "disturbance · 20.9 → 17.9 mm", RED)
    metric_card(slide, Inches(6.72), Inches(1.65), Inches(2.85), "32.7%", "Better final error", "wheel slip · 40.7 → 27.4 mm", GOLD)
    metric_card(slide, Inches(9.72), Inches(1.65), Inches(2.85), "95.1%", "CNN accuracy", "realistic cluttered maps", NAVY)
    rect(slide, Inches(0.72), Inches(3.65), Inches(11.85), Inches(2.15), LIGHT, line=LINE)
    text(slide, Inches(1.05), Inches(3.98), Inches(2.0), Inches(0.35), "Near-oracle", 18, NAVY, True, TITLE_FONT)
    text(slide, Inches(1.05), Inches(4.55), Inches(3.3), Inches(0.7),
         "17.9 mm CNN final error\n15.6 mm oracle upper bound", 15, INK)
    text(slide, Inches(5.05), Inches(3.98), Inches(2.4), Inches(0.35), "Honest trade-off", 18, NAVY, True, TITLE_FONT)
    text(slide, Inches(5.05), Inches(4.55), Inches(6.7), Inches(0.7),
         "Under persistent noise, smoother commands widen steady-state error—"
         "exactly as predicted by the Lyapunov boundary-layer analysis.", 14, MUTED)
    takeaway(slide, "CNN adaptation improves the metric that matters most for each uncertainty type.")
    script = (
        "[Timing: 65 seconds]\n\n"
        "These are the four numbers I want the panel to remember. Under noise, CNN adaptation "
        "reduces chattering by thirty-five point seven percent. After an external disturbance, "
        "final error improves by fourteen point four percent. Under wheel slip, final error "
        "improves by thirty-two point seven percent. On realistic maps, the classifier reaches "
        "ninety-five point one percent accuracy. In the disturbance test, CNN final error is "
        "only two point three millimeters above the oracle. We also report the limitation "
        "honestly: reducing noise-induced chattering increases steady-state error because of "
        "the wider boundary layer."
    )
    note(slide, script); scripts.append(("12. Headline results", script))

    # 13 — Trajectory
    slide = prs.slides.add_slide(blank); header(slide, "Combined Scenario: Recovery After the Push", "Results", 13)
    add_image_contain(slide, FIG / "fig03_trajectory_comparison.png",
                      Inches(0.7), Inches(1.42), Inches(8.65), Inches(4.9))
    rect(slide, Inches(9.62), Inches(1.55), Inches(2.95), Inches(4.45), LIGHT, line=LINE)
    text(slide, Inches(9.92), Inches(1.9), Inches(2.35), Inches(0.38), "Read the graph", 18, NAVY, True, TITLE_FONT)
    bullets(slide, Inches(9.9), Inches(2.55), Inches(2.3), Inches(2.8), [
        "Reference path: dashed",
        "Push at x ≈ 2.4 m",
        "Classical: largest overshoot",
        "CNN: tighter recovery",
        "RL: intermediate response",
    ], 12, spacing=7)
    takeaway(slide, "Adaptive controllers recover more tightly after the same impulsive disturbance.")
    script = (
        "[Timing: 55 seconds]\n\n"
        "This figure shows the actual robot trajectories in the combined scenario. The dashed "
        "line is the desired path. Up to the push, the controllers remain close together. At "
        "approximately two point four meters, the lateral disturbance creates the sharp "
        "deviation. Classical SMC produces the largest overshoot and wider recovery loop. "
        "CNN-adaptive SMC returns more tightly because the disturbance preset increases heading "
        "correction. PPO lies between the two. The key point is that all three controllers "
        "receive the exact same disturbance; the difference comes from parameter adaptation."
    )
    note(slide, script); scripts.append(("13. Trajectory results", script))

    # 14 — Error
    slide = prs.slides.add_slide(blank); header(slide, "Combined Scenario: Error Through Time", "Results", 14)
    add_image_contain(slide, FIG / "fig04_tracking_error.png",
                      Inches(0.7), Inches(1.42), Inches(9.1), Inches(4.92))
    rect(slide, Inches(10.05), Inches(1.58), Inches(2.52), Inches(4.35), WHITE, line=LINE)
    text(slide, Inches(10.32), Inches(1.92), Inches(2.0), Inches(0.35), "Events", 18, NAVY, True, TITLE_FONT)
    for y, tm, event, c in [
        (2.65, "0–8 s", "noise only", TEAL),
        (3.42, "8 s", "external push", RED),
        (4.19, "10–14 s", "wheel slip", GOLD),
        (4.96, "14–20 s", "recovery", PURPLE),
    ]:
        text(slide, Inches(10.3), Inches(y), Inches(0.75), Inches(0.3), tm, 11, c, True)
        text(slide, Inches(11.08), Inches(y), Inches(1.2), Inches(0.3), event, 11, MUTED)
    takeaway(slide, "The time series separates disturbance response, slip response and final recovery.")
    script = (
        "[Timing: 55 seconds]\n\n"
        "The time-series view makes the experiment easier to interpret. During the first eight "
        "seconds, only measurement noise is active. At eight seconds, the external push causes "
        "an immediate error spike. From ten to fourteen seconds, wheel slip slows the robot and "
        "creates a second error increase. After fourteen seconds, all controllers recover. "
        "Classical SMC reaches the highest peak and remains above the adaptive methods for much "
        "of the recovery. This plot confirms that the trajectory difference is not a visual "
        "artifact; it corresponds to lower tracking error over time."
    )
    note(slide, script); scripts.append(("14. Tracking error", script))

    # 15 — Accuracy/smoothness tradeoff
    slide = prs.slides.add_slide(blank); header(slide, "Accuracy vs. Smoothness: The Real Trade-off", "Results", 15)
    add_image_contain(slide, FIG / "fig05_rmse_bars.png",
                      Inches(0.65), Inches(1.45), Inches(5.95), Inches(4.5))
    add_image_contain(slide, FIG / "fig06_chattering_bars.png",
                      Inches(6.75), Inches(1.45), Inches(5.95), Inches(4.5))
    text(slide, Inches(0.9), Inches(5.78), Inches(5.45), Inches(0.4),
         "RMSE stays similar; combined case remains hardest.", 12, MUTED, align=PP_ALIGN.CENTER)
    text(slide, Inches(7.0), Inches(5.78), Inches(5.45), Inches(0.4),
         "CNN strongly reduces chattering under noise and combined.", 12, MUTED, align=PP_ALIGN.CENTER)
    takeaway(slide, "The proposed controller trades a small RMSE change for much smoother actuation.")
    script = (
        "[Timing: 60 seconds]\n\n"
        "These two charts must be read together. On the left, RMSE is similar between classical "
        "and CNN-adaptive control, and the combined scenario remains the hardest. Under noise, "
        "the adaptive controller accepts a small loss in precision. On the right, we see why: "
        "chattering falls dramatically under noise and combined uncertainty. This is not a "
        "free improvement; it is a deliberate engineering trade-off. Smoother angular commands "
        "protect actuators and reduce instability, while the controller preserves competitive "
        "tracking accuracy."
    )
    note(slide, script); scripts.append(("15. Trade-off", script))

    # 16 — Heatmap
    slide = prs.slides.add_slide(blank); header(slide, "Where CNN Adaptation Helps—and Where It Does Not", "Discussion", 16)
    add_image_contain(slide, FIG / "fig09_improvement_heatmap.png",
                      Inches(0.7), Inches(1.42), Inches(7.7), Inches(4.95))
    rect(slide, Inches(8.72), Inches(1.55), Inches(3.85), Inches(4.55), LIGHT, line=LINE)
    text(slide, Inches(9.03), Inches(1.9), Inches(3.2), Inches(0.35), "Interpretation", 18, NAVY, True, TITLE_FONT)
    bullets(slide, Inches(9.0), Inches(2.5), Inches(3.15), Inches(2.95), [
        "Green = adaptive better",
        "Noise: major chattering win",
        "Disturbance/slip: final error win",
        "Red cells expose costs",
        "No claim of universal superiority",
    ], 13)
    takeaway(slide, "Adaptation is scenario-dependent—not uniformly better on every metric.")
    script = (
        "[Timing: 60 seconds]\n\n"
        "The heatmap gives the most honest summary. Green means CNN-adaptive is better; red "
        "means worse. Under noise, the main benefit is smoothness, while final error increases. "
        "Under disturbance and slip, the main benefit is final recovery, while effort or "
        "chattering may increase because the controller becomes more aggressive. Under combined "
        "uncertainty, chattering improves strongly but RMSE changes very little. Therefore, our "
        "claim is not that CNN adaptation wins every cell. Our claim is that it makes the "
        "trade-off context-aware and explainable."
    )
    note(slide, script); scripts.append(("16. Heatmap discussion", script))

    # 17 — Benchmark
    slide = prs.slides.add_slide(blank); header(slide, "Benchmark Against Fuzzy, Oracle and PPO", "Discussion", 17)
    add_image_contain(slide, FIG / "fig10_multicontroller.png",
                      Inches(0.7), Inches(1.42), Inches(8.8), Inches(4.95))
    rect(slide, Inches(9.8), Inches(1.58), Inches(2.75), Inches(4.35), WHITE, line=LINE)
    text(slide, Inches(10.05), Inches(1.93), Inches(2.2), Inches(0.35), "What it proves", 18, NAVY, True, TITLE_FONT)
    bullets(slide, Inches(10.02), Inches(2.55), Inches(2.15), Inches(2.75), [
        "CNN ≈ oracle",
        "CNN beats fuzzy on key recovery metrics",
        "RL best on combined chatter",
        "RL weaker on slip",
    ], 12)
    takeaway(slide, "CNN supervision approaches oracle performance with lower deployment complexity.")
    script = (
        "[Timing: 55 seconds]\n\n"
        "This benchmark provides context beyond the classical baseline. Fuzzy gain scheduling "
        "reduces noise chattering modestly, but CNN and oracle are clearly lower. Under "
        "disturbance, CNN reaches seventeen point nine millimeters compared with the oracle's "
        "fifteen point six. Under slip, CNN and oracle are nearly identical. PPO is best on "
        "combined chattering but weak under slip and disturbance. Overall, the CNN supervisor "
        "offers the strongest balance of performance, interpretability and low runtime cost."
    )
    note(slide, script); scripts.append(("17. Multi-controller benchmark", script))

    # 18 — Digital twin
    slide = prs.slides.add_slide(blank); header(slide, "From Research Code to a Live Digital Twin", "Demonstration", 18)
    add_image_contain(slide, FIG / "fig08_digital_twin.png",
                      Inches(0.65), Inches(1.45), Inches(8.05), Inches(4.95))
    rect(slide, Inches(9.0), Inches(1.55), Inches(3.58), Inches(4.65), NAVY_DARK)
    text(slide, Inches(9.32), Inches(1.9), Inches(2.95), Inches(0.4), "Demo sequence", 20, WHITE, True, TITLE_FONT)
    demo = [
        ("1", "Press T", "Controller Tour"),
        ("2", "Compare", "3-controller benchmark"),
        ("3", "Dual view", "Classical vs CNN paths"),
        ("4", "Replay", "Saved simulation"),
        ("5", "Export", "Metrics as CSV"),
    ]
    for i, (n, a, d) in enumerate(demo):
        y = Inches(2.62 + i * 0.62)
        text(slide, Inches(9.32), y, Inches(0.35), Inches(0.28), n, 13, TEAL, True)
        text(slide, Inches(9.75), y, Inches(0.85), Inches(0.28), a, 12, WHITE, True)
        text(slide, Inches(10.62), y, Inches(1.5), Inches(0.28), d, 10, RGBColor(196, 214, 223))
    text(slide, Inches(9.32), Inches(5.85), Inches(2.9), Inches(0.25),
         "localhost:8000", 12, TEAL, True, align=PP_ALIGN.CENTER)
    takeaway(slide, "The twin makes controller behaviour visible, comparable and replayable.")
    script = (
        "[Timing: 30 seconds, then switch to browser for 3–4 minutes]\n\n"
        "This is the live digital twin. The left panel controls controller mode, scenario, "
        "trajectory and speed. The center visualizes the robot and its path in a hospital "
        "corridor. The right panel shows live metrics. I will now switch to the browser. "
        "First I press T to run the controller tour. Then I open Compare to benchmark all "
        "three controllers. If time allows, I show the dual path view and replay. "
        "\n\nDEMO RECOVERY LINE: If the live demo fails, say: 'The system is also recorded here; "
        "this screenshot and the following quantitative results were generated by the same backend.'"
    )
    note(slide, script); scripts.append(("18. Live demo", script))

    # 19 — Limitations
    slide = prs.slides.add_slide(blank); header(slide, "Limitations: What We Have Not Yet Proven", "Critical reflection", 19)
    limitations = [
        ("Synthetic perception", "Most maps are engineered; realistic set is procedural, not real LiDAR."),
        ("Pre-mission inference", "Current CNN classifies once, not continuously during mission."),
        ("Hand-tuned presets", "CNN chooses parameters but does not learn them end-to-end."),
        ("Simulation scope", "2D kinematics, straight reference and fixed seeds."),
        ("RL maturity", "PPO needs more training and reward shaping before deployment."),
    ]
    for i, (t, d) in enumerate(limitations):
        y = Inches(1.48 + i * 0.94)
        text(slide, Inches(0.82), y + Inches(0.1), Inches(0.45), Inches(0.35), f"{i+1}", 16, RED, True, TITLE_FONT)
        text(slide, Inches(1.45), y, Inches(2.7), Inches(0.35), t, 15, NAVY, True)
        text(slide, Inches(4.25), y, Inches(7.8), Inches(0.55), d, 13, MUTED)
        rect(slide, Inches(1.45), y + Inches(0.66), Inches(10.65), Inches(0.012), LINE, radius=False)
    takeaway(slide, "The results establish feasibility—not physical-robot generalisation.")
    script = (
        "[Timing: 50 seconds]\n\n"
        "We want to be precise about the limits of the evidence. Most occupancy maps are "
        "synthetic, and the realistic set is still procedural rather than real LiDAR. The "
        "CNN currently classifies once before the mission. Parameter presets are hand-tuned, "
        "and the controller is evaluated in a two-dimensional kinematic simulation with fixed "
        "seeds. Finally, PPO remains an experimental extension. Therefore, this work demonstrates "
        "feasibility and a reproducible architecture; it does not yet prove generalisation to "
        "a physical robot or arbitrary terrain."
    )
    note(slide, script); scripts.append(("19. Limitations", script))

    # 20 — Conclusion
    slide = prs.slides.add_slide(blank); header(slide, "Conclusion & Next Steps", "Conclusion", 20)
    rect(slide, Inches(0.72), Inches(1.52), Inches(5.8), Inches(4.72), NAVY_DARK)
    text(slide, Inches(1.05), Inches(1.88), Inches(5.1), Inches(0.4), "What we demonstrated", 21, WHITE, True, TITLE_FONT)
    bullets(slide, Inches(1.05), Inches(2.56), Inches(4.95), Inches(2.75), [
        "Context-aware SMC adaptation",
        "35.7% chattering reduction",
        "Up to 32.7% recovery improvement",
        "95.1% realistic-map classification",
        "Operational 3D digital twin",
    ], 14, RGBColor(225, 235, 240))
    rect(slide, Inches(6.82), Inches(1.52), Inches(5.8), Inches(4.72), LIGHT, line=LINE)
    text(slide, Inches(7.15), Inches(1.88), Inches(5.1), Inches(0.4), "What comes next", 21, NAVY, True, TITLE_FONT)
    bullets(slide, Inches(7.15), Inches(2.56), Inches(4.95), Inches(2.75), [
        "Real LiDAR occupancy maps",
        "Continuous online inference",
        "Monte-Carlo confidence intervals",
        "Longer PPO curriculum training",
        "Physical robot validation",
    ], 14)
    takeaway(slide, "Explainable AI can improve robust control without replacing the control law.")
    script = (
        "[Timing: 60 seconds]\n\n"
        "To conclude, we built a complete perception-control-demonstration pipeline. The CNN "
        "makes sliding mode control context-aware without replacing its interpretable control "
        "law. The main evidence is a thirty-five point seven percent chattering reduction under "
        "noise, up to thirty-two point seven percent better final recovery, ninety-five point "
        "one percent classification on realistic maps, and a working digital twin. The next "
        "step is to move from procedural maps and fixed scenarios to real LiDAR, continuous "
        "online inference, Monte Carlo evaluation and a physical differential-drive robot. "
        "The core message is that explainable AI can improve robust control while preserving "
        "engineering transparency."
    )
    note(slide, script); scripts.append(("20. Conclusion", script))

    # 21 — Questions
    slide = prs.slides.add_slide(blank); set_bg(slide, NAVY_DARK)
    rect(slide, Inches(0), Inches(0), W, Inches(0.14), TEAL, radius=False)
    text(slide, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.6),
         "Thank you", 42, WHITE, True, TITLE_FONT, align=PP_ALIGN.CENTER)
    text(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.5),
         "Questions & discussion", 22, TEAL, True, align=PP_ALIGN.CENTER)
    rect(slide, Inches(2.1), Inches(3.0), Inches(9.1), Inches(1.55), RGBColor(15, 75, 105), line=TEAL)
    text(slide, Inches(2.4), Inches(3.35), Inches(8.5), Inches(0.32),
         "CNN-Adaptive SMC · PPO · Digital Twin", 20, WHITE, True, align=PP_ALIGN.CENTER)
    text(slide, Inches(2.4), Inches(3.92), Inches(8.5), Inches(0.3),
         "github.com/LikhitaYerra/smc_cnn", 13, RGBColor(205, 224, 232), align=PP_ALIGN.CENTER)
    add_image_contain(slide, ASSETS / "aivancity_logo.png",
                      Inches(5.23), Inches(5.35), Inches(2.87), Inches(0.85))
    text(slide, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.3),
         "Supervisor · Prof. Vishvjit Thakar", 12, WHITE, align=PP_ALIGN.CENTER)
    script = (
        "[Timing: Q&A]\n\n"
        "Thank you for your attention. We are ready for your questions.\n\n"
        "LIKELY QUESTION — Why CNN if random forest is slightly better on realistic maps?\n"
        "ANSWER: The CNN integrates naturally with image features and scales to richer spatial "
        "inputs. The random forest result confirms that the current cues are separable; the CNN "
        "is the extensible architecture for future LiDAR and camera maps.\n\n"
        "LIKELY QUESTION — Is 100% accuracy overfitting?\n"
        "ANSWER: It reflects engineered class separability, not real-world performance. That is "
        "why we separately report 95.1% on cluttered maps and state physical validation as future work.\n\n"
        "LIKELY QUESTION — Why is CNN not better on every metric?\n"
        "ANSWER: Smoothing and tracking precision conflict. The Lyapunov bound predicts that a "
        "wider boundary layer reduces chattering while widening steady-state error.\n\n"
        "LIKELY QUESTION — Why not use PPO alone?\n"
        "ANSWER: Current PPO is competitive on combined chattering but weaker under slip and "
        "disturbance. CNN supervision is more reliable, interpretable and cheaper at runtime."
    )
    note(slide, script); scripts.append(("21. Q&A", script))

    # Backup 22 — preset table
    slide = prs.slides.add_slide(blank); header(slide, "Backup: Scenario-Specific SMC Presets", "Appendix", 22)
    table_data = [
        ["Parameter", "Normal", "Noise", "Disturb.", "Slip", "Combined"],
        ["λx", "2.00", "2.00", "2.15", "2.20", "2.15"],
        ["λy", "2.00", "2.10", "2.30", "2.20", "2.35"],
        ["λθ", "1.00", "0.95", "1.05", "1.00", "1.00"],
        ["kv", "0.30", "0.30", "0.33", "0.35", "0.35"],
        ["kω", "0.80", "0.78", "0.88", "0.86", "0.88"],
        ["φ", "0.50", "0.58", "0.52", "0.62", "0.62"],
        ["αω", "0.950", "0.965", "0.945", "0.960", "0.965"],
    ]
    shape = slide.shapes.add_table(len(table_data), len(table_data[0]),
                                   Inches(1.0), Inches(1.55), Inches(11.35), Inches(4.75))
    table = shape.table
    for r, row in enumerate(table_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (LIGHT if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.name = BODY_FONT
                p.font.size = Pt(13)
                p.font.bold = r == 0 or c == 0
                p.font.color.rgb = WHITE if r == 0 else INK
    takeaway(slide, "Preset values encode the engineering rationale shown on Slide 9.")
    script = (
        "[Backup slide — use only if asked about exact controller parameters]\n\n"
        "These are the complete scenario-specific parameter values. Point out that noise "
        "raises phi and smoothing, disturbance raises lateral and angular correction, and "
        "slip raises both position gains. The combined preset is a compromise."
    )
    note(slide, script); scripts.append(("22. Backup presets", script))

    # Backup 23 — stability
    slide = prs.slides.add_slide(blank); header(slide, "Backup: Stability Argument", "Appendix", 23)
    rect(slide, Inches(0.85), Inches(1.6), Inches(11.65), Inches(3.7), LIGHT, line=LINE)
    text(slide, Inches(1.2), Inches(1.95), Inches(10.9), Inches(0.55),
         "V = ½(sₓ² + sᵧ² + sθ²) ≥ 0", 25, NAVY, True, TITLE_FONT, align=PP_ALIGN.CENTER)
    text(slide, Inches(1.2), Inches(2.88), Inches(10.9), Inches(0.6),
         "V̇ ≤ −kω Σ sᵢ tanh(sᵢ/φ) + d̄ ‖s‖", 23, NAVY_DARK, False, TITLE_FONT, align=PP_ALIGN.CENTER)
    text(slide, Inches(1.2), Inches(3.8), Inches(10.9), Inches(0.75),
         "V̇ < 0 outside a bounded region  →  sliding surfaces are uniformly ultimately bounded",
         17, TEAL, True, align=PP_ALIGN.CENTER)
    text(slide, Inches(1.2), Inches(4.7), Inches(10.9), Inches(0.4),
         "Wider φ → lower chattering, larger ultimate tracking-error bound", 15, RED, True, align=PP_ALIGN.CENTER)
    takeaway(slide, "The theory predicts the exact smoothness–accuracy trade-off observed in Results.")
    script = (
        "[Backup slide — use if asked about stability]\n\n"
        "We use the composite Lyapunov function based on the three sliding surfaces. Because "
        "x times tanh of x over phi is positive away from zero, the derivative is negative "
        "outside a bounded disturbance-dependent region. This gives uniform ultimate boundedness. "
        "Inside the boundary layer, the error bound grows with phi. That is why noise smoothing "
        "reduces chattering but can increase steady-state error."
    )
    note(slide, script); scripts.append(("23. Backup stability", script))

    # Backup 24 — technical stack
    slide = prs.slides.add_slide(blank); header(slide, "Backup: Implementation & Reproducibility", "Appendix", 24)
    stack = [
        ("Simulation", "Python · NumPy · Δt = 0.01 s", NAVY),
        ("AI", "PyTorch · EnvironmentCNN · PPO", TEAL),
        ("Backend", "FastAPI · WebSocket · REST export", PURPLE),
        ("Frontend", "React · Three.js · Vite", GOLD),
        ("Artifacts", "models · CSV · replay · figures", RED),
    ]
    for i, (t, d, c) in enumerate(stack):
        y = Inches(1.5 + i * 0.93)
        rect(slide, Inches(1.0), y, Inches(11.25), Inches(0.7), WHITE, line=LINE)
        rect(slide, Inches(1.0), y, Inches(0.09), Inches(0.7), c, radius=False)
        text(slide, Inches(1.3), y + Inches(0.16), Inches(2.0), Inches(0.3), t, 14, c, True)
        text(slide, Inches(3.65), y + Inches(0.16), Inches(7.8), Inches(0.3), d, 14, INK)
    takeaway(slide, "One command launches the trained models, simulation API and digital twin.")
    script = (
        "[Backup slide — use if asked about implementation]\n\n"
        "The simulation and evaluation are Python-based. PyTorch handles CNN and PPO. FastAPI "
        "streams state over WebSocket, while React and Three.js render the 3D twin. Runs can be "
        "recorded, replayed and exported. The project includes trained checkpoints and scripts "
        "for regeneration."
    )
    note(slide, script); scripts.append(("24. Backup implementation", script))

    prs.save(str(OUT))

    # Standalone script is useful if Presenter View is unavailable.
    lines = [
        "# AI Clinic Defense — Speaker Script",
        "",
        "Core deck: Slides 1–21. Slides 22–24 are backup/Q&A only.",
        "",
    ]
    for title_value, script in scripts:
        lines.extend([f"## {title_value}", "", script, ""])
    SCRIPT_OUT.write_text("\n".join(lines), encoding="utf-8")

    return OUT, len(prs.slides), scripts


if __name__ == "__main__":
    out, count, scripts = build()
    print(f"Presentation: {out}")
    print(f"Slides: {count}")
    print(f"Notes sections: {len(scripts)}")
    print(f"Speaker script: {SCRIPT_OUT}")
